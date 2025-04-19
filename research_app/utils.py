import os
import time

import fitz  # PyMuPDF
from django.conf import settings
from docx import Document as DocxDocument # Avoid confusion with Django Document
from docx.shared import Inches
from google import genai
from pptx import Presentation

# -- Global Variables --

GEMINI_MODEL = "gemini-2.0-flash"
GEMINI_CLIENT = genai.Client(api_key=settings.GEMINI_API_KEY)

# --- Text Extraction ---

def extract_text_from_pdf(file_path):
    """Extracts text from a PDF file."""
    try:
        doc = fitz.open(file_path)
        text = ""
        metadata = doc.metadata
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            page_text = page.get_text("text")
             # Simple attempt to get chapter/section titles (heuristic)
            blocks = page.get_text("blocks")
            current_title = f"Page {page_num + 1}" # Default context
            # Look for larger font sizes maybe indicating titles - highly heuristic
            # A more robust solution would involve analyzing font styles, layout etc.
            # This is a basic placeholder for context.
            # for b in blocks:
            #     # You might analyze b[4] (text), b[5] (font size info is complex)
            #     pass # Add logic here if needed
            text += f"\n--- {current_title} ---\n{page_text}\n"
        return text, metadata
    except Exception as e:
        print(f"Error extracting PDF {os.path.basename(file_path)}: {e}")
        return None, None

def extract_text_from_docx(file_path):
    """Extracts text and attempts to identify headings from a DOCX file."""
    try:
        doc = DocxDocument(file_path)
        text = ""
        current_heading = "General Content"
        for para in doc.paragraphs:
            # Check if the paragraph style suggests it's a heading
            # This relies on standard heading styles (Heading 1, Heading 2, etc.)
            if para.style and para.style.name.startswith('Heading'):
                current_heading = para.text.strip()
                if current_heading: # Add heading marker only if not empty
                     text += f"\n--- Section: {current_heading} ---\n"
            elif para.text.strip(): # Add paragraph text if not empty
                 text += f"{para.text}\n"
        return text, {"title": os.path.basename(file_path)} # Basic metadata
    except Exception as e:
        print(f"Error extracting DOCX {os.path.basename(file_path)}: {e}")
        return None, None

def extract_text_from_pptx(file_path):
    """Extracts text from a PPTX file, including slide titles."""
    try:
        prs = Presentation(file_path)
        text = ""
        for i, slide in enumerate(prs.slides):
            slide_title = f"Slide {i + 1}"
            try:
                if slide.shapes.title:
                    slide_title = f"Slide {i + 1}: {slide.shapes.title.text.strip()}"
            except AttributeError:
                pass # No title shape

            text += f"\n--- {slide_title} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    # Append text from text frames, ignoring empty ones
                    shape_text = shape.text.strip()
                    if shape_text:
                       text += f"{shape_text}\n"
            # Extract text from notes slide if present
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    text += f"\n--- Notes for {slide_title} ---\n{notes_text}\n"

        return text, {"title": os.path.basename(file_path)} # Basic metadata
    except Exception as e:
        print(f"Error extracting PPTX {os.path.basename(file_path)}: {e}")
        return None, None

def extract_text_from_txt(file_path):
    """Extracts text from a TXT file."""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        return text, {"title": os.path.basename(file_path)}
    except Exception as e:
        print(f"Error reading TXT {os.path.basename(file_path)}: {e}")
        return None, None

def extract_text(document_obj):
    """Main text extraction dispatcher."""
    file_path = document_obj.file.path
    _, extension = os.path.splitext(file_path)
    extension = extension.lower()

    text = None
    metadata = None # Placeholder for potential metadata extraction

    print(f"Attempting extraction for: {document_obj.original_filename}")
    if extension == '.pdf':
        text, metadata = extract_text_from_pdf(file_path)
    elif extension == '.docx':
        text, metadata = extract_text_from_docx(file_path)
    elif extension == '.pptx':
        text, metadata = extract_text_from_pptx(file_path)
    elif extension == '.txt':
        text, metadata = extract_text_from_txt(file_path)
    else:
        print(f"Unsupported file type: {extension}")
        document_obj.processing_log = f"Unsupported file type: {extension}"
        document_obj.status = 'error'
        document_obj.save()
        return None # Indicate failure

    if text:
        document_obj.extracted_text = text # Save extracted text if desired (can be large)
        document_obj.status = 'converted'
        print(f"Extraction successful for: {document_obj.original_filename}")
    else:
        document_obj.processing_log = f"Failed to extract text from {document_obj.original_filename}"
        document_obj.status = 'error'
        print(f"Extraction failed for: {document_obj.original_filename}")

    document_obj.save()
    return text, metadata # Return text for immediate use

# --- Gemini Interaction ---

def query_gemini_single_doc(text, query, filename):
    """Queries Gemini model for an answer within a single document's text."""
    if not GEMINI_MODEL:
        return "Error: Gemini model not configured.", ""
    if not text or not text.strip():
         return "Document contains no extractable text.", ""

    prompt = f"""
    Źródło dokumentu: {filename}

    Przeanalizuj następujący tekst dokumentu wyłącznie na podstawie podanego tekstu.
    Odpowiedz na pytanie: "{query}"

    Instrukcje:
    1. Podaj krótką odpowiedź na pytanie na podstawie tekstu.
    2. Jeśli dokument zawiera informacje, które mogą pomóc w odpowiedzi, uwzględnij 1-3 bezpośrednie cytaty z tekstu do wsparcia swojej odpowiedzi. Formatuj cytaty tak: "tekst cytatu...".
    3. Uwzględnij kontekst dla cytatów, takie jak numery stron lub tytuły sekcji, jeśli są one obecne w tekście (np. z '--- Strona X ---' lub '--- Sekcja: Y ---' markerów).
    4. Jeśli tekst dokumentu *nie* zawiera informacji dotyczących pytania, wyraźnie i uczciwie stwierdź: "Ten dokument nie zawiera informacji dotyczących pytania: '{query}'." Nie wymyślaj informacji ani nie wyciągaj wniosków poza podany tekst.
    5. Struktura swojej odpowiedzi powinna być klarowna. Zacznij od bezpośredniej odpowiedzi, po której następują wspominające cytaty (jeśli są), lub stwierdzenie, że nie znaleziono informacji dotyczących pytania.

    Tekst dokumentu:
    --- POCZĄTEK TEKSTU ---
    {text[:1_000_000]}
    --- KONIEC TEKSTU ---

    Twoja odpowiedź:
    """

    # Simple retry mechanism
    max_retries = 2
    for attempt in range(max_retries):
        try:
            response = GEMINI_CLIENT.models.generate_content(
                model=GEMINI_MODEL,
                contents=prompt,
            )
            answer_text = response.text.strip()

            # Extract quotes for structured data (optional, mainly for summary prompt)
            quotes = [q.strip() for q in answer_text.split('"') if q.strip() and len(q) > 5] # Simple extraction
            return answer_text, quotes

        except Exception as e:
            print(f"Gemini API error (Attempt {attempt + 1}/{max_retries}) on {filename}: {e}")
            if "quota" in str(e).lower() or "rate limit" in str(e).lower():
                 time.sleep(5 * (attempt + 1)) # Exponential backoff for rate limits
            elif attempt == max_retries - 1:
                return f"Error: Failed to get response from LLM after {max_retries} attempts. Last error: {e}", []
            time.sleep(2) # General delay between retries
    return f"Error: Failed to get response from LLM after {max_retries} attempts.", []


def query_gemini_summary(all_answers_text, query):
    """Generates a summary answer based on findings from all documents."""
    if not GEMINI_MODEL:
        return "Error: Gemini model not configured."

    prompt = f"""
    Opracuj syntetyczną odpowiedź na pytanie: "{query}"

    Podstaw swoją syntezę wyłącznie na podstawie podanych wyników, które zostały wyciągnięte z różnych dokumentów.

    Instrukcje:
    1. Utwórz zwięzłą syntezę, która bezpośrednio odpowiada na pytanie.
    2. Uwzględnij kluczowe punkty z analiz poszczególnych dokumentów.
    3. Podczas cytowania informacji, odwołuj się do dokumentu, który został jawnie wymieniony w wynikach (np. "Zgodnie z 'raport.pdf'...", lub "Zgodnie z 'prezentacja.pptx', ...").
    4. Jeśli wiele dokumentów dostarcza sprzecznych informacji, uwzględnij te sprzeczności.
    5. Jeśli wyniki wskazują, że żaden z dokumentów nie zawierał informacji dotyczących pytania, wyraźnie stwierdź to.
    6. Nie dodawaj informacji, które nie są obecne w podanych wynikach.

    Wyniki z dokumentów:
    --- POCZĄTEK WYNIKÓW ---
    {all_answers_text[:100000]}
    --- KONIEC WYNIKÓW ---

    Twoja odpowiedź:
    """
    # Limit combined text size for safety

    max_retries = 2
    for attempt in range(max_retries):
         try:
            response = GEMINI_CLIENT.models.generate_content(
                model=GEMINI_MODEL,
                contents=prompt,
            )
            return response.text.strip()
         except Exception as e:
            print(f"Gemini API error during summary (Attempt {attempt + 1}/{max_retries}): {e}")
            if "quota" in str(e).lower() or "rate limit" in str(e).lower():
                 time.sleep(5 * (attempt + 1))
            elif attempt == max_retries - 1:
                return f"Error: Failed to get summary response from LLM after {max_retries} attempts. Last error: {e}"
            time.sleep(2)
    return f"Error: Failed to get summary response from LLM after {max_retries} attempts."


# --- Docx Generation ---

def initialize_report(query):
    """Creates a new docx document and adds initial title and query."""
    doc = DocxDocument()
    doc.add_heading(f'Raport dotyczący pytania: "{query}"', level=0)
    doc.add_paragraph(query)
    doc.add_paragraph() # Add some space
    return doc

def add_answer_to_report(doc, filename, answer_text):
    """Adds the analysis results for a single document to the report."""
    doc.add_heading(f'Analiza: {filename}', level=2)
    # Add answer paragraph - handle potential None or empty answer
    answer_paragraph = answer_text if answer_text else "Nie wygenerowano odpowiedzi lub wystąpił błąd."
    doc.add_paragraph(answer_paragraph)
     # Add a visual separator
    doc.add_paragraph("---")
    doc.add_paragraph() # Add some space

def add_summary_to_report(doc, summary_text):
    """Adds the final summary section to the report."""
    doc.add_heading('Odpowiedź syntetyczna', level=1)
    doc.add_paragraph(summary_text if summary_text else "Nie można wygenerować syntezy.")

def save_report(doc, session):
    """Saves the docx document to the media/reports directory."""
    reports_dir = os.path.join(settings.MEDIA_ROOT, 'reports')
    os.makedirs(reports_dir, exist_ok=True)

    # Sanitize query for filename (simple example)
    safe_query_part = "".join(c if c.isalnum() else "_" for c in session.query[:30])
    filename = f"report_{session.session_id}_{safe_query_part}.docx"
    filepath = os.path.join(reports_dir, filename)

    try:
        doc.save(filepath)
        session.report_filename = filename # Store only the filename
        session.save()
        print(f"Report saved successfully: {filepath}")
        return filepath
    except Exception as e:
        print(f"Error saving report {filename}: {e}")
        session.status = 'failed'
        session.error_message = f"Error saving report: {e}"
        session.save()
        return None
